import pandas as pd
from pptx import Presentation
from io import BytesIO
import streamlit as st
import traceback

def export_to_excel(parsed_content):
    try:
        if isinstance(parsed_content, pd.DataFrame):
            df = parsed_content
        else:
            # Assuming the parsed content is a string representation of a table
            # Split the content into rows and columns
            rows = []
            for line in parsed_content.split('\n'):
                if line:
                    rows.append(line.split('\t'))
            
            # Create DataFrame with rows and columns
            df = pd.DataFrame(rows[1:], columns=rows[0])
        
        excel_buffer = BytesIO()
        with pd.ExcelWriter(excel_buffer, engine='openpyxl') as writer:
            df.to_excel(writer, index=False, sheet_name='Parsed Content')
            worksheet = writer.sheets['Parsed Content']
            for col in worksheet.columns:
                max_length = 0
                column = col[0].column_letter
                for cell in col:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(cell.value)
                    except:
                        pass
                adjusted_width = (max_length + 2)
                worksheet.column_dimensions[column].width = adjusted_width
        excel_buffer.seek(0)
        
        st.download_button(
            label="📊 Export to Excel",
            data=excel_buffer,
            file_name="parsed_content.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )
    except Exception as e:
        st.error(f"Error exporting to Excel: {str(e)}")
        st.code(traceback.format_exc())



def export_to_powerpoint(parsed_content):
    try:
        prs = Presentation()
        
        # Add title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        title.text = "Web Scraping Results"
        subtitle.text = "Extracted Content"
        
        # Add content slide
        content_slide = prs.slides.add_slide(prs.slide_layouts[1])
        content_slide.shapes.title.text = "Parsed Content"
        content_shape = content_slide.placeholders[1]
        content_shape.text = parsed_content
        
        # Save presentation
        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        st.download_button(
            label="📑 Export to PowerPoint",
            data=pptx_buffer,
            file_name="parsed_content.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    except Exception as e:
        st.error(f"Error exporting to PowerPoint: {str(e)}")
        st.code(traceback.format_exc())
